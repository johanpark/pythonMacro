from pptx import Presentation
from openpyxl import load_workbook
from collections import defaultdict
import pandas as pd


def extract_words_from_excel(excel_file):
    words = []
    wb = load_workbook(excel_file)
    ws = wb.active
    for row in ws.iter_rows(values_only=True):
        for cell in row:
            if isinstance(cell, str):
                words.append(cell)
    return words


def count_words_in_pptx(pptx_file, words):
    prs = Presentation(pptx_file)
    pptx_text = " ".join([shape.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame])

    word_counter = defaultdict(int)

    for word in words:
        word_count = pptx_text.lower().count(word.lower())  # 대소문자 구분 없이 단어 빈도 계산
        word_counter[word] += word_count

        # 입력한 단어와 비슷한 대소문자 조합도 빈도 계산
        similar_words = [w for w in pptx_text.split() if w.lower() == word.lower()]
        word_counter[word] += len(similar_words)

    return word_counter


# 경로 설정
excel_file = '/Users/it1454/Downloads/타겟 엑셀.xlsx'
pptx_file = '/Users/it1454/Downloads/NICE그룹 성과관리 시스템 구축_제안서.pptx'

excel_words = extract_words_from_excel(excel_file)
word_counts = count_words_in_pptx(pptx_file, excel_words)

word_count_data = [{'Word': word, 'Count': count} for word, count in word_counts.items()]
df = pd.DataFrame(word_count_data)

# 엑셀 파일로 저장
output_excel_file = '결과_파일.xlsx'
df.to_excel(output_excel_file, index=False)

print(f"결과가 '{output_excel_file}'에 저장되었습니다.")
